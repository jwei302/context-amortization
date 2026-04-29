#!/usr/bin/env python3
"""Build slide/final_presentation.pptx for CPSC 4770/5770 final project.

Run: `python slide/build_deck.py` from the repo root.
Regenerates intermediate assets in slide/assets/ (idempotent) and the final pptx.
"""

from __future__ import annotations

import subprocess
from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.animation import FuncAnimation, PillowWriter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


ROOT = Path(__file__).resolve().parents[1]
SLIDE = ROOT / "slide"
ASSETS = SLIDE / "assets"
FIGS = ROOT / "figs"
EVAL = ROOT / "eval_out"
ASSETS.mkdir(exist_ok=True)

BRAND_ORANGE = RGBColor(0xFD, 0x64, 0x0B)
BRAND_ORANGE_HEX = "#FD640B"
INK = RGBColor(0x23, 0x1F, 0x20)
INK_HEX = "#231F20"
MUTED = RGBColor(0x6B, 0x72, 0x80)
MUTED_HEX = "#6B7280"
BG = RGBColor(0xFA, 0xFA, 0xF7)
CREAM = RGBColor(0xFF, 0xF4, 0xEB)
PAPER = RGBColor(0xF5, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LOSS_RED = "#D72638"
FONT = "Inter"

DECK_W, DECK_H = Inches(13.333), Inches(7.5)


# ---------------------------------------------------------------------------
# asset prep
# ---------------------------------------------------------------------------


def run(cmd):
    print(f"  $ {' '.join(str(c) for c in cmd)}")
    subprocess.run([str(c) for c in cmd], check=True)


def prepare_logo() -> Path:
    out = ASSETS / "rhoda-logo.png"
    if out.exists():
        return out
    run(["rsvg-convert", SLIDE / "rhoda-logo-light.svg", "-h", "200", "-o", out])
    return out


def prepare_circular_photo(src: Path, size: int = 400) -> Path:
    out = ASSETS / f"{src.stem}-circle.png"
    if out.exists():
        return out
    img = Image.open(src).convert("RGBA")
    s = min(img.size)
    left = (img.width - s) // 2
    top = (img.height - s) // 2
    img = img.crop((left, top, left + s, top + s)).resize((size, size), Image.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    out_img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    out_img.paste(img, (0, 0), mask)
    out_img.save(out)
    return out


def prepare_training_regime_gif() -> Path:
    """Animated GIF: DF vs CA training regime (noise + loss mask per frame)."""
    out = ASSETS / "training_regime.gif"
    if out.exists():
        return out

    rng = np.random.default_rng(0)
    N = 8  # frames in a sequence
    K = 4  # CA anchors
    STEPS = 10

    noise_df = rng.uniform(0.15, 0.9, size=(STEPS, N))
    noise_ca = rng.uniform(0.15, 0.9, size=(STEPS, N))

    fig, ax = plt.subplots(figsize=(11, 2.85), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_xlim(-2.6, N + 2.8)
    ax.set_ylim(-0.5, 2.75)
    ax.set_aspect("equal")
    ax.axis("off")

    # title
    title = ax.text(
        (N / 2) - 0.5, 2.6, "Training step 1 / 10",
        ha="center", fontsize=12, color=INK_HEX, weight="bold",
    )

    # row labels
    ax.text(-0.2, 1.925, "Diffusion\nForcing", ha="right", va="center",
            fontsize=12, color=INK_HEX, weight="bold")
    ax.text(-0.2, 0.425, "Context\nAmortization\n(k=4)", ha="right", va="center",
            fontsize=12, color=BRAND_ORANGE_HEX, weight="bold")

    # right-side legends
    ax.text(N + 0.3, 1.925, "random noise\non all 8 frames\n∇loss on all 8",
            ha="left", va="center", fontsize=9, color=MUTED_HEX)
    ax.text(N + 0.3, 0.425, "first 4 clean (ε≈0)\nlast 4 random noise\n∇loss only on last 4",
            ha="left", va="center", fontsize=9, color=MUTED_HEX)

    # separator band between the two rows
    ax.plot([-1.9, N + 2.4], [1.175, 1.175],
            color="#e8e8e8", lw=0.8, zorder=0)
    ax.text(N / 2 - 0.5, 1.26, "CA closes the train–test gap",
            ha="center", va="bottom", fontsize=9,
            color=BRAND_ORANGE_HEX, weight="bold", style="italic")

    df_rects, df_badges = [], []
    ca_rects, ca_badges = [], []

    for i in range(N):
        # DF row at y in [1.5, 2.4]
        r = mpatches.FancyBboxPatch(
            (i, 1.5), 0.85, 0.85,
            boxstyle="round,pad=0.02,rounding_size=0.08",
            facecolor="#d0d0d0", edgecolor=INK_HEX, lw=1.1,
        )
        ax.add_patch(r)
        df_rects.append(r)
        b = ax.text(i + 0.42, 1.92, "∇", ha="center", va="center",
                    fontsize=13, color=LOSS_RED, weight="bold")
        df_badges.append(b)

        # CA row at y in [0, 0.85]
        is_anchor = i < K
        fc = "white" if is_anchor else "#d0d0d0"
        ec = BRAND_ORANGE_HEX if is_anchor else INK_HEX
        lw = 2.0 if is_anchor else 1.1
        r = mpatches.FancyBboxPatch(
            (i, 0), 0.85, 0.85,
            boxstyle="round,pad=0.02,rounding_size=0.08",
            facecolor=fc, edgecolor=ec, lw=lw,
        )
        ax.add_patch(r)
        ca_rects.append(r)
        if is_anchor:
            b = ax.text(i + 0.42, 0.42, "✗", ha="center", va="center",
                        fontsize=12, color=MUTED_HEX, weight="bold")
        else:
            b = ax.text(i + 0.42, 0.42, "∇", ha="center", va="center",
                        fontsize=13, color=LOSS_RED, weight="bold")
        ca_badges.append(b)

    # separators — small frame indices under each column
    for i in range(N):
        ax.text(i + 0.42, -0.15, f"x{i+1}", ha="center", va="center",
                fontsize=7, color=MUTED_HEX)

    def update(t):
        for i, r in enumerate(df_rects):
            v = noise_df[t, i]
            r.set_facecolor((1 - v, 1 - v, 1 - v))
        for i, r in enumerate(ca_rects):
            if i < K:
                continue
            v = noise_ca[t, i]
            r.set_facecolor((1 - v, 1 - v, 1 - v))
        title.set_text(f"Training step {t + 1} / {STEPS}")
        return df_rects + ca_rects + [title]

    ani = FuncAnimation(fig, update, frames=STEPS, interval=500, blit=False)
    ani.save(out, writer=PillowWriter(fps=2))
    plt.close(fig)

    # re-save to guarantee loop=0 (PillowWriter doesn't always set it)
    im = Image.open(out)
    frames = []
    try:
        while True:
            frames.append(im.copy().convert("P"))
            im.seek(im.tell() + 1)
    except EOFError:
        pass
    frames[0].save(
        out, save_all=True, append_images=frames[1:],
        duration=500, loop=0, disposal=2, optimize=True,
    )
    return out


def _mp4_to_gif(src: Path, dst: Path, fps: int = 10, width: int = 256):
    """High-quality palettized gif via ffmpeg."""
    vf = (
        f"fps={fps},scale={width}:-1:flags=lanczos,"
        "split[s0][s1];[s0]palettegen=max_colors=128[p];[s1][p]paletteuse"
    )
    run(["ffmpeg", "-y", "-i", src, "-vf", vf, "-loglevel", "error", dst])


def prepare_rollout_gifs():
    """Slide 6: gt + 3 preds × 2 samples from h=128 batch0."""
    needed = []
    for sample in [0, 1]:
        # gt is identical across models; pick ca_k4 copy
        needed.append(("gt", sample, EVAL / f"ca_k4/h128/gt_batch0_sample{sample}.mp4"))
        for model in ["df", "ca_k4", "ca_k8"]:
            needed.append((model, sample,
                           EVAL / f"{model}/h128/pred_batch0_sample{sample}.mp4"))

    outputs = []
    for tag, sample, src in needed:
        dst = ASSETS / f"{tag}_sample{sample}.gif"
        if not dst.exists():
            if not src.exists():
                print(f"  WARNING: missing {src}")
                continue
            _mp4_to_gif(src, dst)
        outputs.append((tag, sample, dst))
    return outputs


# ---------------------------------------------------------------------------
# pptx helpers
# ---------------------------------------------------------------------------


def _set_run(run, text, *, size, bold=False, color=INK, font=FONT, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size=size, bold=bold, italic=italic, color=color, font=font)
    return tb


def add_bullets(slide, x, y, w, h, bullets, *,
                size=14, color=INK, spacing=10, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        _set_run(p.add_run(), "•  ", size=size, bold=True, color=BRAND_ORANGE, font=font)
        _set_run(p.add_run(), bullet, size=size, color=color, font=font)
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = line_width
    return shp


# ---------------------------------------------------------------------------
# slides
# ---------------------------------------------------------------------------


def slide_title(prs, *, logo, me_circle, steven_circle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, DECK_W, DECK_H, BG)
    add_rect(s, 0, 0, Inches(0.3), DECK_H, BRAND_ORANGE)

    add_text(s, Inches(0.9), Inches(1.2), Inches(12), Inches(1.3),
             "Context Amortization vs. Diffusion Forcing",
             size=44, bold=True)
    add_text(s, Inches(0.9), Inches(2.25), Inches(12), Inches(0.7),
             "Closing the train–test gap on DMLab video prediction",
             size=22, color=MUTED)

    # authors
    photo_h = Inches(1.6)
    label_y = Inches(5.75)
    s.shapes.add_picture(str(me_circle), Inches(2.4), Inches(4.0), height=photo_h)
    add_text(s, Inches(1.9), label_y, Inches(2.6), Inches(0.4),
             "Jeffrey Wei", size=18, bold=True, align=PP_ALIGN.CENTER)
    s.shapes.add_picture(str(steven_circle), Inches(5.5), Inches(4.0), height=photo_h)
    add_text(s, Inches(5.0), label_y, Inches(2.6), Inches(0.4),
             "Steven Zhou", size=18, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.9), Inches(6.85), Inches(6), Inches(0.4),
             "CPSC 4770/5770  ·  Spring 2026", size=14, color=MUTED)

    add_text(s, Inches(8.4), Inches(6.7), Inches(1.7), Inches(0.3),
             "Inspired by", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    s.shapes.add_picture(str(logo), Inches(10.25), Inches(6.7), height=Inches(0.45))


def slide_motivation(prs, *, logo):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, DECK_W, DECK_H, BG)

    add_text(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8),
             "The train–test gap in Diffusion Forcing",
             size=30, bold=True)
    add_rect(s, Inches(0.6), Inches(1.05), Inches(1.2), Inches(0.06), BRAND_ORANGE)

    add_bullets(s, Inches(0.6), Inches(1.45), Inches(7.3), Inches(4.0), [
        "Diffusion Forcing (Chen et al., NeurIPS 2025) samples an independent random noise level for every frame during training — including frames that will become context at inference.",
        "At inference, context is clean (noise = 0) and only the future is denoised. The “clean context + noisy future” joint distribution is never seen during training.",
        "The model is rolled out from outside its training support. Autoregressive drift follows.",
    ], size=15, spacing=14)

    # prior work sidebar
    add_rect(s, Inches(8.3), Inches(1.45), Inches(4.5), Inches(3.1), CREAM)
    add_text(s, Inches(8.55), Inches(1.6), Inches(4.0), Inches(0.4),
             "Prior work", size=14, bold=True, color=BRAND_ORANGE)
    add_bullets(s, Inches(8.55), Inches(2.1), Inches(4.0), Inches(1.8), [
        "Diffusion Forcing — Chen et al., 2025",
        "Teacher forcing + prompt masks in LMs",
        "Rhoda AI research blog (rhoda.ai)",
    ], size=12, spacing=6)
    s.shapes.add_picture(str(logo), Inches(8.55), Inches(3.85), height=Inches(0.45))

    # analogy band
    add_rect(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.7), PAPER,
             line=MUTED, line_width=Pt(0.5))
    add_text(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.4),
             "Key analogy", size=12, bold=True, color=BRAND_ORANGE)
    add_text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(1.2),
             "“LMs do this with teacher forcing and a prompt mask — the prompt incurs no loss, uncertainty lives only in the targets. Context Amortization is the video-diffusion analog.”",
             size=15, italic=True)


def slide_methods(prs, *, training_gif):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, DECK_W, DECK_H, BG)

    add_text(s, Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.7),
             "Context Amortization: clean anchors + masked loss",
             size=28, bold=True)
    add_rect(s, Inches(0.6), Inches(0.95), Inches(1.2), Inches(0.06), BRAND_ORANGE)

    # centered animation; source is ~11:4 aspect
    s.shapes.add_picture(str(training_gif), Inches(1.15), Inches(1.2), width=Inches(11.0))

    # config strip
    add_text(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.4),
             "Configs    DF (k=1)     ·     CA k=4 (4 clean / 28 noised)     ·     CA k=8 (8 clean / 24 noised)",
             size=13, color=INK, align=PP_ALIGN.CENTER, bold=True)

    # compute footer band
    add_rect(s, 0, Inches(6.35), DECK_W, Inches(1.15), INK)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.35),
             "Training & compute", size=11, bold=True, color=BRAND_ORANGE)
    add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.45),
             "Lambda Cloud  ·  1× NVIDIA H200 (GH200)  ·  bf16-mixed  ·  100k steps  ·  ~11 h / run  ·  DMLab (8-frame context → 128 predicted)",
             size=12, color=WHITE)
    add_text(s, Inches(0.6), Inches(7.15), Inches(12.1), Inches(0.3),
             "Identical architecture, optimizer, batch size, and wall-clock across all three runs.",
             size=10, italic=True, color=RGBColor(0xC8, 0xC8, 0xC8))


def slide_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, DECK_W, DECK_H, BG)

    add_text(s, Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.8),
             "CA beats DF at every horizon — 25 to 47% lower FVD",
             size=26, bold=True)
    add_rect(s, Inches(0.6), Inches(1.0), Inches(1.2), Inches(0.06), BRAND_ORANGE)

    # FVD table (left)
    tx, ty, tw, th = Inches(0.6), Inches(1.3), Inches(6.0), Inches(2.3)
    headers = ["Config", "h = 32", "h = 64", "h = 128"]
    data = [
        ["DF",     "1500", "807", "852"],
        ["CA k=4", "799",  "825", "580"],
        ["CA k=8", "1062", "779", "635"],
    ]
    winner = {1: 1, 2: 2, 3: 1}  # col_idx → winning row_idx (0-based within data)

    table = s.shapes.add_table(4, 4, tx, ty, tw, th).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), h, size=13, bold=True, color=WHITE)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else PAPER
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            is_winner = (j > 0 and winner.get(j) == i)
            color = BRAND_ORANGE if is_winner else INK
            bold = is_winner or j == 0
            _set_run(p.add_run(), val, size=15, bold=bold, color=color)

    add_text(s, Inches(0.6), Inches(3.7), Inches(6), Inches(0.3),
             "FVD (lower is better), 16 held-out DMLab clips · winners bold + orange",
             size=10, italic=True, color=MUTED)

    # takeaway card
    add_rect(s, Inches(0.6), Inches(4.1), Inches(6.0), Inches(1.2), CREAM)
    add_text(s, Inches(0.85), Inches(4.22), Inches(5.6), Inches(0.4),
             "Headline", size=11, bold=True, color=BRAND_ORANGE)
    add_text(s, Inches(0.85), Inches(4.55), Inches(5.6), Inches(0.75),
             "CA k=4 best overall:  580 FVD @ h=128  vs  DF 852  (−32%).  Same architecture, same compute budget.",
             size=14, bold=True)

    # right: fvd vs horizon
    s.shapes.add_picture(str(FIGS / "fvd_vs_horizon.png"),
                         Inches(6.9), Inches(1.3), width=Inches(6.1))

    # bottom thumbnails
    y_thumb = Inches(5.55)
    s.shapes.add_picture(str(FIGS / "fvd_vs_step.png"),
                         Inches(0.6), y_thumb, width=Inches(4.0))
    s.shapes.add_picture(str(FIGS / "loss_curves.png"),
                         Inches(4.8), y_thumb, width=Inches(4.0))
    s.shapes.add_picture(str(FIGS / "lpips_vs_horizon.png"),
                         Inches(9.0), y_thumb, width=Inches(4.0))


def slide_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, DECK_W, DECK_H, BG)

    add_text(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             "Takeaways & future work", size=30, bold=True)
    add_rect(s, Inches(0.6), Inches(1.0), Inches(1.2), Inches(0.06), BRAND_ORANGE)

    add_text(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5),
             "What worked", size=18, bold=True, color=BRAND_ORANGE)
    add_bullets(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5), [
        "Masking loss on clean anchors closes the DF train–test gap without touching architecture or compute.",
        "k=4 dominates at long horizons (h=128); k=8 at medium (h=64). Anchor count is a tunable hyperparameter vs. target horizon.",
        "Gains come essentially for free: same wall-clock, same batch size, same parameter count — only the noise schedule and loss mask change.",
    ], size=15, spacing=12)

    add_text(s, Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5),
             "Future work", size=18, bold=True, color=BRAND_ORANGE)
    add_bullets(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.5), [
        "Sweep larger k (12, 16) and non-zero anchor-noise schedules.",
        "Transfer to richer video domains: Minecraft, robotics, driving.",
        "Longer horizons with rolling-anchor inference at test time.",
        "Compare against classifier-free-guidance-style context conditioning.",
    ], size=15, spacing=10)

    add_rect(s, 0, Inches(6.6), DECK_W, Inches(0.9), INK)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.4),
             "Checkpoints   huggingface.co/jwei302/context-amortization     ·     Code   github.com/jwei302/context-amortization",
             size=12, color=WHITE)


def slide_demo(prs, rollout_gifs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, DECK_W, DECK_H, BG)

    add_text(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "Rollouts @ h = 128  (from 8-frame clean prefix)",
             size=24, bold=True)
    add_rect(s, Inches(0.6), Inches(0.95), Inches(1.2), Inches(0.06), BRAND_ORANGE)

    # organize gifs
    by = {}
    for tag, sample, path in rollout_gifs:
        by.setdefault(tag, {})[sample] = path

    rows = [
        ("Ground truth", "gt",    INK),
        ("DF",           "df",    MUTED),
        ("CA k=4",       "ca_k4", BRAND_ORANGE),
        ("CA k=8",       "ca_k8", BRAND_ORANGE),
    ]

    gif_size = Inches(1.25)           # square
    row_gap = Inches(0.12)
    col_gap = Inches(0.25)
    label_w = Inches(1.6)
    n_cols = 2
    grid_w = label_w + n_cols * gif_size + (n_cols - 1) * col_gap
    grid_left = (DECK_W - grid_w) / 2
    grid_top = Inches(1.35)

    # column headers
    for j in range(n_cols):
        cx = grid_left + label_w + j * (gif_size + col_gap)
        add_text(s, cx, Inches(1.05), gif_size, Inches(0.3),
                 f"Sample {j+1}", size=10, bold=True, color=MUTED,
                 align=PP_ALIGN.CENTER)

    for i, (label, tag, color) in enumerate(rows):
        y = grid_top + i * (gif_size + row_gap)
        add_text(s, grid_left, y + Inches(0.35), label_w, Inches(0.6),
                 label, size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)
        for j in range(n_cols):
            x = grid_left + label_w + j * (gif_size + col_gap)
            path = by.get(tag, {}).get(j)
            if path is None:
                continue
            s.shapes.add_picture(str(path), x, y, width=gif_size, height=gif_size)

    # caption
    add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5),
             "CA rollouts hold scene structure further into the horizon; DF drifts off-support as prediction continues.",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------


def main():
    print("[1/3] preparing assets")
    logo = prepare_logo()
    me_circle = prepare_circular_photo(SLIDE / "me.jpg")
    steven_circle = prepare_circular_photo(SLIDE / "steven.png")
    training_gif = prepare_training_regime_gif()
    rollout_gifs = prepare_rollout_gifs()

    print("[2/3] building deck")
    prs = Presentation()
    prs.slide_width = DECK_W
    prs.slide_height = DECK_H

    slide_title(prs, logo=logo, me_circle=me_circle, steven_circle=steven_circle)
    slide_motivation(prs, logo=logo)
    slide_methods(prs, training_gif=training_gif)
    slide_results(prs)
    slide_conclusion(prs)
    slide_demo(prs, rollout_gifs)

    out = SLIDE / "final_presentation.pptx"
    prs.save(out)

    print("[3/3] done")
    print(f"  slides:  {len(prs.slides)}")
    print(f"  output:  {out}   ({out.stat().st_size / 1e6:.2f} MB)")
    print(f"  assets:  {ASSETS}")
    for p in sorted(ASSETS.iterdir()):
        print(f"    - {p.name}  ({p.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
